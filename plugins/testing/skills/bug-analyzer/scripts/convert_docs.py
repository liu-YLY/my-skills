"""
文档格式转换工具（降级方案）：将 .docx / .xlsx / .pptx 转换为 AI 可读格式。

首选方案为 Microsoft MarkItDown（pip install markitdown[docx,pptx,xlsx,xls]），
支持 PDF、.xls 等更多格式且质量更高。本脚本用于 MarkItDown 不可用时降级。

用法：
    python scripts/convert_docs.py <文件路径>              # 转换单个文件
    python scripts/convert_docs.py <目录路径>              # 转换目录下所有支持的文件
    python scripts/convert_docs.py <目录路径> --recursive  # 递归转换子目录

输出：在源文件同目录下生成同名的 .md 或 .csv 文件。

依赖安装：
    pip install python-docx openpyxl python-pptx

注意：不支持 .xls 和 PDF，请优先使用 MarkItDown。
"""

import argparse
import re
import sys
from datetime import date, datetime
from pathlib import Path

SUPPORTED_EXTENSIONS = {".docx", ".xlsx", ".pptx"}


def _escape_table_cell(val: str) -> str:
    return str(val).replace("|", "\\|").replace("\n", "<br>")


def _has_markdown_list_prefix(text: str) -> bool:
    return bool(re.match(r"^[-*+]\s", text))


def convert_docx_to_md(file_path: Path) -> Path:
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(file_path)
    lines = []

    for para in doc.paragraphs:
        style_name = para.style.name.lower() if para.style else ""
        text = para.text.strip()
        if not text:
            lines.append("")
            continue

        is_list = False
        numPr = para._element.find(qn("w:pPr"))
        if numPr is not None:
            numPr_elem = numPr.find(qn("w:numPr"))
            is_list = numPr_elem is not None

        if is_list or _has_markdown_list_prefix(text):
            clean = re.sub(r"^[-*+]\s*", "", text, count=1)
            lines.append(f"- {clean}")
        elif "heading 1" in style_name:
            lines.append(f"# {text}")
        elif "heading 2" in style_name:
            lines.append(f"## {text}")
        elif "heading 3" in style_name:
            lines.append(f"### {text}")
        elif "heading" in style_name:
            lines.append(f"#### {text}")
        else:
            lines.append(text)

    for table in doc.tables:
        lines.append("")
        for i, row in enumerate(table.rows):
            cells = [_escape_table_cell(cell.text.strip()) for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
        lines.append("")

    content = "\n".join(lines)
    content = re.sub(r"\n{3,}", "\n\n", content)

    output = file_path.with_suffix(".md")
    output.write_text(content, encoding="utf-8")
    return output


def convert_xlsx_to_csv(file_path: Path) -> list[Path]:
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    outputs = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            if all(cell is None or str(cell).strip() == "" for cell in row):
                continue
            cells = []
            for cell in row:
                if cell is None:
                    cells.append("")
                elif isinstance(cell, datetime):
                    cells.append(cell.isoformat())
                elif isinstance(cell, date):
                    cells.append(cell.isoformat())
                elif isinstance(cell, float):
                    if cell == int(cell) and abs(cell) < 1e15:
                        cells.append(str(int(cell)))
                    else:
                        cells.append(repr(cell))
                else:
                    val = str(cell)
                    if "," in val or '"' in val or "\n" in val:
                        val = '"' + val.replace('"', '""') + '"'
                    cells.append(val)
            rows.append(",".join(cells))

        if not rows:
            continue

        if len(wb.sheetnames) == 1:
            output = file_path.with_suffix(".csv")
        else:
            safe_name = sheet_name.replace("/", "_").replace("\\", "_")
            output = file_path.with_name(f"{file_path.stem}_{safe_name}.csv")

        output.write_text("\n".join(rows), encoding="utf-8")
        outputs.append(output)

    wb.close()
    return outputs


def convert_pptx_to_md(file_path: Path) -> Path:
    from pptx import Presentation

    prs = Presentation(file_path)
    lines = []

    def _extract_shapes(shapes):
        for shape in shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                table = shape.table
                for ri, row in enumerate(table.rows):
                    cells = [
                        _escape_table_cell(cell.text.strip()) for cell in row.cells
                    ]
                    lines.append("| " + " | ".join(cells) + " |")
                    if ri == 0:
                        lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
            if shape.shape_type == 6:
                _extract_shapes(shape.shapes)

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        _extract_shapes(slide.shapes)
        lines.append("")

    content = "\n".join(lines)
    content = re.sub(r"\n{3,}", "\n\n", content)

    output = file_path.with_suffix(".md")
    output.write_text(content, encoding="utf-8")
    return output


def convert_file(file_path: Path) -> list[Path]:
    ext = file_path.suffix.lower()
    if ext == ".docx":
        return [convert_docx_to_md(file_path)]
    elif ext == ".xlsx":
        return convert_xlsx_to_csv(file_path)
    elif ext == ".pptx":
        return [convert_pptx_to_md(file_path)]
    else:
        return []


def collect_files(path: Path, recursive: bool) -> list[Path]:
    if path.is_file():
        return [path] if path.suffix.lower() in SUPPORTED_EXTENSIONS else []

    pattern = "**/*" if recursive else "*"
    return [f for f in path.glob(pattern) if f.suffix.lower() in SUPPORTED_EXTENSIONS]


def main():
    parser = argparse.ArgumentParser(description="转换 Word/Excel/PPT 为 AI 可读格式")
    parser.add_argument("path", help="文件或目录路径")
    parser.add_argument("--recursive", "-r", action="store_true", help="递归处理子目录")
    args = parser.parse_args()

    target = Path(args.path)
    if not target.exists():
        print(f"错误：路径不存在 {target}", file=sys.stderr)
        sys.exit(1)

    files = collect_files(target, args.recursive)
    if not files:
        print(f"未找到支持的文件（{', '.join(SUPPORTED_EXTENSIONS)}）")
        sys.exit(0)

    success, failed = 0, 0
    for f in files:
        try:
            outputs = convert_file(f)
            for out in outputs:
                print(f"  ✓ {f.name} → {out.name}")
            success += 1
        except Exception as e:
            print(f"  ✗ {f.name} → 失败: {e}", file=sys.stderr)
            failed += 1

    print(f"\n完成：{success} 成功, {failed} 失败")


if __name__ == "__main__":
    main()
